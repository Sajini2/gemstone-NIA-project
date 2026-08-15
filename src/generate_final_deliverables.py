import os
import csv
from docx import Document
from docx.shared import Inches, Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt

def generate_final_report():
    doc = Document()
    
    # Title Page
    doc.add_heading('Horizon Campus', 0).alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_heading('Module: IT41033', 1).alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_heading('Gemstone Price Prediction using Nature-Inspired Feature Selection', 1).alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    doc.add_paragraph('\n\n')
    p = doc.add_paragraph('Group Members:')
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p = doc.add_paragraph('Sajini - Index Number: [Sajini Index]\nBuddhika Janadari - Index Number: [Buddhika Index]')
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph('\n')
    p = doc.add_paragraph('Module Lecturer: Mr. Sanka Wijewardene')
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_page_break()

    # Abstract
    doc.add_heading('Abstract', 1)
    doc.add_paragraph("This report investigates the application of Nature-Inspired Algorithms (NIA), specifically Genetic Algorithms (GA) and Particle Swarm Optimization (PSO), for feature selection in predicting gemstone prices. Using the Diamonds and Cubic Zirconia datasets, we demonstrate that metaheuristic optimization significantly reduces feature dimensionality (up to 55.6%) while maintaining high predictive accuracy (over 0.97 R²).")
    
    # Introduction
    doc.add_heading('1. Introduction', 1)
    doc.add_paragraph("Gemstone pricing is a complex task dependent on multiple correlated features, traditionally the '4 Cs' (carat, cut, color, clarity) along with spatial dimensions. This project aims to apply machine learning models (Random Forest and XGBoost) to predict prices. To mitigate the curse of dimensionality and enhance interpretability, we employ Genetic Algorithms (GA) and Particle Swarm Optimization (PSO) to autonomously select optimal feature subsets.")
    
    # Literature Review
    doc.add_heading('2. Literature Review', 1)
    doc.add_paragraph("Previous studies highlight the effectiveness of tree-based ensembles like Random Forest and XGBoost for tabular regression tasks. Concurrently, metaheuristic algorithms like GA and PSO have been widely adopted for feature selection due to their ability to efficiently navigate vast combinatorial search spaces, avoiding local optima that plague traditional greedy search methods.")

    # Methodology
    doc.add_heading('3. Methodology', 1)
    
    doc.add_heading('3.1 Genetic Algorithm (GA) Feature Selection', 2)
    doc.add_paragraph("Sajini utilized the Diamonds dataset, applying a Random Forest regressor as the fitness evaluator. The GA encodes feature subsets as binary chromosomes. The fitness function balances the cross-validated R² score against a penalty for the number of features. Standard crossover and mutation operators were employed across 40 generations with a population of 40.")
    
    doc.add_heading('3.2 Particle Swarm Optimization (PSO) Feature Selection', 2)
    doc.add_paragraph("Buddhika utilized the Cubic Zirconia (Gemstone) dataset, using an XGBoost regressor. The PSO algorithm encodes particle positions as continuous velocities transformed into a binary mask via a sigmoid function. The swarm consisted of 30 particles over 40 iterations. The fitness function precisely matched the GA penalty: Fitness = Average CV R² - (0.001 * N_features). Cognitive and social parameters (c1, c2) were set to 1.5, with inertia weight decaying from 0.9 to 0.4.")

    # Results
    doc.add_heading('4. Results', 1)
    
    csv_path = 'results/final_comparison.csv'
    if os.path.exists(csv_path):
        with open(csv_path, 'r') as f:
            data = list(csv.reader(f))
        table = doc.add_table(rows=len(data), cols=len(data[0]))
        table.style = 'Table Grid'
        for i, row in enumerate(data):
            for j, val in enumerate(row):
                cell = table.cell(i, j)
                cell.text = str(val)
                if i == 0:
                    for paragraph in cell.paragraphs:
                        for run in paragraph.runs:
                            run.font.bold = True
    
    # Add charts
    charts = [
        ('results/comparison_metrics_chart.png', "Figure 1: Grouped bar chart comparing RMSE, MAE, and R²."),
        ('results/ga_vs_pso_convergence.png', "Figure 2: Convergence comparison of GA and PSO."),
        ('results/complexity_comparison.png', "Figure 3: Model complexity (features) vs Training time.")
    ]
    for path, caption in charts:
        if os.path.exists(path):
            doc.add_picture(path, width=Inches(5.5))
            doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER
            cap = doc.add_paragraph(caption)
            cap.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # Discussion
    doc.add_heading('5. Discussion', 1)
    doc.add_paragraph("The results confirm our hypothesis. GA and PSO successfully isolated the '4 Cs' while dropping redundant spatial dimensions (x, y, z), reducing dimensionality by 33.3% and 55.6% respectively. Accuracy drops were negligible (-0.108% for GA and -0.159% for PSO), proving that NIA feature selection provides significantly simpler and domain-intuitive models. Limitations include the use of separate datasets, making direct algorithm comparisons difficult.")

    # Conclusion
    doc.add_heading('6. Conclusion', 1)
    doc.add_paragraph("This project successfully demonstrated the power of Genetic Algorithms and Particle Swarm Optimization in feature selection for gemstone price prediction. By reducing feature counts by up to 55.6% with under 0.2% accuracy loss, we established the efficiency of these metaheuristics. Future work involves testing on a shared single dataset, developing a hybrid GA-PSO model, and comparing against Deep Learning architectures.")

    # References
    doc.add_heading('7. References', 1)
    refs = [
        "[1] S. Agrawal, 'Diamonds Dataset,' Kaggle, 2017.",
        "[2] 'Cubic Zirconia Dataset,' Kaggle.",
        "[3] F. Pedregosa et al., 'Scikit-learn: Machine Learning in Python,' JMLR, 2011.",
        "[4] T. Chen and C. Guestrin, 'XGBoost: A Scalable Tree Boosting System,' KDD, 2016.",
        "[5] J. Kennedy and R. Eberhart, 'Particle swarm optimization,' ICNN, 1995."
    ]
    for r in refs:
        doc.add_paragraph(r)

    doc.save('report/Final_Report.docx')

def generate_presentation():
    prs = Presentation()
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Gemstone Price Prediction using NIA"
    slide.placeholders[1].text = "Sajini & Buddhika Janadari\nModule: IT41033 | Horizon Campus\nLecturer: Mr. Sanka Wijewardene"
    slide.notes_slide.notes_text_frame.text = "[Both] Welcome and introduction."

    # Slide 2: Problem Statement
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Problem Statement & Research Question"
    slide.placeholders[1].text = "- Gemstone pricing depends on highly correlated features.\n- Can Nature-Inspired Algorithms (GA/PSO) autonomously select feature subsets that reduce complexity while maintaining predictive accuracy?"
    slide.notes_slide.notes_text_frame.text = "[Sajini] Explain the complexity of gemstone pricing and introduce the research question."

    # Slide 3: Datasets
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Datasets Used"
    slide.placeholders[1].text = "- Diamonds Dataset (Sajini): 53,940 instances, 10 features.\n- Cubic Zirconia Dataset (Buddhika): 26,967 instances, 10 features.\n- Target: Price. Core features: Carat, Cut, Color, Clarity, Depth, Table, X, Y, Z."
    slide.notes_slide.notes_text_frame.text = "[Sajini] Describe the data sources and preprocessing steps."

    # Slide 4: Methodology
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Methodology Overview"
    slide.placeholders[1].text = "Data Preprocessing \u2192 Baseline Models (RF/XGB) \u2192 Metaheuristic Optimization (GA/PSO) \u2192 Final Models \u2192 Evaluation"
    slide.notes_slide.notes_text_frame.text = "[Buddhika] Walk through the pipeline."

    # Slide 5: GA Approach
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Genetic Algorithm (GA) Approach"
    slide.placeholders[1].text = "- Presenter: Sajini\n- Encoding: Binary Chromosome\n- Fitness: Avg CV R2 - 0.001*N_features\n- Operators: Tournament Selection, Uniform Crossover, Bit-flip Mutation"
    slide.notes_slide.notes_text_frame.text = "[Sajini] Detail the GA implementation choices."

    # Slide 6: PSO Approach
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Particle Swarm Optimization (PSO) Approach"
    slide.placeholders[1].text = "- Presenter: Buddhika\n- Encoding: Continuous velocities transformed via Sigmoid to binary.\n- Fitness: Same penalty formula.\n- Update: Standard velocity updates with decaying inertia."
    slide.notes_slide.notes_text_frame.text = "[Buddhika] Explain the continuous to binary PSO transformation."

    # Slide 7: Convergence
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Convergence Comparison"
    if os.path.exists('results/ga_vs_pso_convergence.png'):
        slide.shapes.add_picture('results/ga_vs_pso_convergence.png', PptxInches(1.5), PptxInches(2), height=PptxInches(4.5))
    slide.notes_slide.notes_text_frame.text = "[Buddhika] Both algorithms successfully plateaued at highly optimal fitness levels rapidly."

    # Slide 8: Results & Metrics
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Results Comparison"
    if os.path.exists('results/comparison_metrics_chart.png'):
        slide.shapes.add_picture('results/comparison_metrics_chart.png', PptxInches(1.5), PptxInches(2), height=PptxInches(4.5))
    slide.notes_slide.notes_text_frame.text = "[Sajini] Compare the optimized models to the full-feature baselines."

    # Slide 9: Key Findings
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Key Findings & Hypothesis Support"
    slide.placeholders[1].text = "- Both GA and PSO selected the '4 Cs', dropping redundant spatial features.\n- GA reduced features by 33.3%, PSO by 55.6%.\n- Accuracy (R2) dropped by less than 0.2% in both cases.\n- Hypothesis Supported: Complexity heavily reduced with virtually no accuracy penalty."
    slide.notes_slide.notes_text_frame.text = "[Buddhika] State the statistical improvements."

    # Slide 10: Conclusion & Future Work
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Conclusion & Future Work"
    slide.placeholders[1].text = "- Conclusion: NIA effectively isolates critical predictive features, providing simpler, domain-intuitive models.\n- Future Work: Shared single dataset for direct comparison, Hybrid GA-PSO architecture, Deep Learning baselines."
    slide.notes_slide.notes_text_frame.text = "[Both] Wrap up presentation."

    prs.save('report/Final_Presentation.pptx')

def generate_contributions():
    doc = Document()
    doc.add_heading('Individual Contribution Statement', 0)
    
    doc.add_heading('Sajini', 1)
    doc.add_paragraph("- Phase 1: Diamonds Dataset EDA and Preprocessing.")
    doc.add_paragraph("- Phase 2: Random Forest Baseline model implementation.")
    doc.add_paragraph("- Phase 3: Genetic Algorithm (GA) feature selection development.")
    doc.add_paragraph("- Phase 4: Final evaluation of GA-selected subsets.")
    doc.add_paragraph("- Report: Drafted Abstract, Introduction, Literature Review, and GA Methodology sections.")
    
    doc.add_heading('Buddhika Janadari', 1)
    doc.add_paragraph("- Phase 1: Cubic Zirconia (Gemstone) Dataset EDA and Preprocessing.")
    doc.add_paragraph("- Phase 2: XGBoost Baseline model implementation.")
    doc.add_paragraph("- Phase 3: Particle Swarm Optimization (PSO) feature selection development.")
    doc.add_paragraph("- Phase 4 & 5: Final evaluation scripts, comparison analytics, and master table generation.")
    doc.add_paragraph("- Report & Delivery: Drafted PSO Methodology, Results, Discussion, and Conclusion sections. Compiled final slide deck and merged report.")
    
    doc.save('report/Individual_Contribution_Statement.docx')

if __name__ == '__main__':
    generate_final_report()
    generate_presentation()
    generate_contributions()
    print("All deliverables generated.")
